"""
FastAPI backend bridge for SF Agentic AI.

This API is intentionally thin:
- Keeps existing `main.py` CLI behavior intact.
- Runs jobs in background subprocesses.
- Persists run/upload state under `data/api`.
"""

from __future__ import annotations

import json
import hashlib
import mimetypes
import os
import re
import subprocess
import sys
import threading
import textwrap
import uuid
from contextlib import asynccontextmanager
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, List, Literal, Optional, Tuple

from dotenv import load_dotenv
from fastapi import FastAPI, File, HTTPException, Request, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse, RedirectResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel, Field

from config import settings
from tools.healthcheck import preflight_lm_studio

try:
    from docx import Document  # type: ignore
    from docx.shared import Pt  # type: ignore
except Exception:  # pragma: no cover - optional import guard
    Document = None
    Pt = None

try:
    from openpyxl import Workbook, load_workbook  # type: ignore
    from openpyxl.styles import Alignment, Font  # type: ignore
except Exception:  # pragma: no cover - optional import guard
    Workbook = None
    load_workbook = None
    Alignment = None
    Font = None

try:
    from pptx import Presentation  # type: ignore
    from pptx.util import Inches, Pt as PptPt  # type: ignore
except Exception:  # pragma: no cover - optional import guard
    Presentation = None
    Inches = None
    PptPt = None

load_dotenv()

ROOT_DIR = Path(__file__).resolve().parent.parent
DATA_DIR = ROOT_DIR / "data" / "api"
UPLOAD_DIR = DATA_DIR / "uploads"
LOG_DIR = DATA_DIR / "logs"
RUNS_FILE = DATA_DIR / "runs.json"
UPLOADS_FILE = DATA_DIR / "uploads.json"
UI_CONFIG_FILE = DATA_DIR / "ui_config.json"
OUTPUT_DIR = ROOT_DIR / "data" / "output"
DEFAULT_CONTROL_DIR = ROOT_DIR / "data" / "inbox" / "control"
STATIC_DIR = ROOT_DIR / "backend" / "static"
SCRIPTS_DIR = ROOT_DIR / "scripts"

SUPPORTED_EXTS = {
    ".pdf",
    ".txt",
    ".md",
    ".csv",
    ".tsv",
    ".json",
    ".doc",
    ".docx",
    ".ppt",
    ".pptx",
    ".xls",
    ".xlsx",
    ".png",
    ".jpg",
    ".jpeg",
}
ALLOWED_INTENTS = {
    "memo_only",
    "memo_plus_distribution",
    "memo_plus_holdbook",
    "full_e2e",
    "custom",
    "analysis_brief",
    "ops_sop",
    "meeting_minutes",
    "pipeline_report",
    "market_intel",
    "data_analysis",
}
ALLOWED_ROUTE_MODES = {"auto", "full", "targeted"}
ALLOWED_POST_CREDIT = {"ask", "stop", "holdbook", "distribution", "hybrid"}
ALLOWED_HUMAN_GATE = {"ask", "approve", "revise", "stop"}
ALLOWED_OUTPUT_FORMATS = {"txt", "md", "docx", "xlsx", "pptx", "pdf"}
ALLOWED_CREDIT_DECISION_MODES = {"analysis_only", "decision_recommendation"}
SUPPORTED_CLOUD_API_PROVIDERS = ("perplexity", "gemini", "chatgpt", "claude", "tavily")
DEFAULT_OUTPUT_FORMATS = ["txt", "docx", "xlsx"]
ARTIFACT_SCHEMA_VERSION = 1
DEFAULT_CREDIT_THRESHOLDS = {
    "dscr_min": 1.2,
    "icr_min": 2.0,
    "net_leverage_max": 4.5,
    "current_ratio_min": 1.0,
}
DEFAULT_CONDITION_PRECEDENTS = [
    "Final legal due diligence satisfactory to lender counsel",
    "Insurance certificate and assignments completed",
    "Signed covenant schedule and reporting calendar confirmed",
]
DEFAULT_OUTPUT_TEMPLATES = {
    "credit_memo": (
        "Use sections: Executive Summary, Borrower Profile, Deal Structure, Financial Analysis, "
        "Risk & Mitigants, Covenants, Recommendation. Include a final decision table."
    ),
    "sop": (
        "Use sections: Purpose, Scope, Roles, Inputs, Procedure Steps, Controls, Escalation, "
        "Appendix. Keep action steps numbered."
    ),
    "meeting_minutes": (
        "Use sections: Meeting Metadata, Attendance, Agenda, Key Discussion Points, Decisions, "
        "Action Items (owner + due date), Follow-up."
    ),
    "pipeline_report": (
        "Use sections: Snapshot KPIs, Pipeline by Stage, Movement Since Last Report, Risk Flags, "
        "Resource Needs, Next Actions. Include tabular summary."
    ),
    "analysis_brief": (
        "Use sections: Executive Summary, Key Findings, Evidence, Analysis, Risks, Open Items. "
        "Do not force credit approval language unless requested."
    ),
    "market_intel": (
        "Use sections: Scope, Market Summary, Comparable Signals, Risks, Implications, "
        "Verification Notes. Mark items requiring source verification."
    ),
    "data_analysis": (
        "Use sections: Data Summary, Data Quality, Key Metrics, Trends, Exceptions, "
        "Recommendations. Use tables where possible."
    ),
}

OUTPUT_FAMILY_ALIASES = {
    "credit_memo_v2": "credit_memo",
    "credit_committee_memo": "credit_memo",
    "memo_only": "credit_memo",
    "memo_plus_distribution": "credit_memo",
    "memo_plus_holdbook": "credit_memo",
    "full_e2e": "credit_memo",
    "sop": "ops_sop",
    "standard_sop": "ops_sop",
    "ops_sop_v2": "ops_sop",
    "analysis": "analysis_brief",
    "analysis_brief_v1": "analysis_brief",
    "meeting_minutes_v1": "meeting_minutes",
    "pipeline_report_v1": "pipeline_report",
    "market_intel_brief": "market_intel",
    "data_report": "data_analysis",
    "autonomous": "analysis_brief",
    "not_applicable": "analysis_brief",
}
OUTPUT_FAMILY_KINDS = {
    "credit_memo": "credit_memo",
    "analysis_brief": "analysis_brief",
    "ops_sop": "sop",
    "meeting_minutes": "meeting_minutes",
    "pipeline_report": "pipeline_report",
    "data_analysis": "data_analysis",
    "market_intel": "market_intel",
}


def _normalize_output_family(value: str, fallback: str = "credit_memo") -> str:
    key = str(value or "").strip().lower().replace("-", "_").replace(" ", "_")
    if not key:
        key = fallback
    normalized = OUTPUT_FAMILY_ALIASES.get(key, key)
    return normalized if normalized in OUTPUT_FAMILY_KINDS else fallback


def _output_family_for_request(req: "RunRequest") -> str:
    for raw in (req.report_format, req.output_template_id, req.intent, req.sop_format):
        family = _normalize_output_family(str(raw or ""), fallback="")
        if family:
            return family
    return "credit_memo"
DEFAULT_WORKER_REGISTRY = [
    {"id": "ingestion", "label": "Ingestion", "desc": "Parse and normalize inputs", "enabled": True, "category": "core"},
    {"id": "market_intel", "label": "Market Intel", "desc": "Research and external benchmarking", "enabled": True, "category": "core"},
    {"id": "analysis_parallel", "label": "Analysis", "desc": "Financial and risk analysis", "enabled": True, "category": "core"},
    {"id": "financial_modeler", "label": "Financial Modeler", "desc": "Projections and sensitivity build", "enabled": True, "category": "core"},
    {"id": "compliance", "label": "Compliance", "desc": "Policy and covenant checks", "enabled": True, "category": "core"},
    {"id": "memo_architect", "label": "Memo Architect", "desc": "Synthesis and final memo writing", "enabled": True, "category": "core"},
    {"id": "senior_advisor", "label": "Senior Advisor", "desc": "Quality assurance and decision framing", "enabled": True, "category": "core"},
    {"id": "admin_ops", "label": "Admin Ops", "desc": "Meeting minutes and SOP drafting", "enabled": True, "category": "support"},
    {"id": "data_analyst", "label": "Data Analyst", "desc": "Pipeline tracking, reporting, and data analysis", "enabled": True, "category": "support"},
]
DEFAULT_MODEL_REGISTRY = [
    {
        "id": "qwen/qwen3.5-9b",
        "label": "Qwen 3.5 9B (LM Studio)",
        "provider": "local",
        "llm_provider": "lm_studio",
        "endpoint": "http://127.0.0.1:1234/v1",
        "enabled": True,
        "is_default": True,
    },
    {
        "id": "lmstudio/qwen/qwen3.5-9b",
        "label": "Qwen 3.5 9B (Hermes Gateway)",
        "provider": "local",
        "llm_provider": "hermes",
        "endpoint": "http://127.0.0.1:18789",
        "enabled": False,
        "is_default": False,
    },
    {
        "id": "gpt-4.1",
        "label": "GPT-4.1 (Cloud)",
        "provider": "cloud",
        "llm_provider": "hermes",
        "endpoint": "",
        "enabled": False,
        "is_default": False,
    },
]


def _default_cloud_apis() -> dict:
    return {
        provider: {
            "enabled": False,
            "api_key": "",
            "endpoint": "",
            "model": "",
        }
        for provider in SUPPORTED_CLOUD_API_PROVIDERS
    }

_STORE_LOCK = threading.Lock()
_PROC_LOCK = threading.Lock()
_ACTIVE_PROCS: Dict[str, subprocess.Popen] = {}
_SCHEDULER_STOP = threading.Event()


def _now_iso() -> str:
    return datetime.now(timezone.utc).isoformat()


def _parse_schedule_time(raw: str) -> Optional[datetime]:
    value = str(raw or "").strip()
    if not value:
        return None
    normalized = value[:-1] + "+00:00" if value.endswith("Z") else value
    normalized = re.sub(r"([+-]\d{2})(\d{2})$", r"\1:\2", normalized)
    try:
        parsed = datetime.fromisoformat(normalized)
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=f"scheduled_at must be an ISO datetime: {exc}")
    if parsed.tzinfo is None:
        parsed = parsed.replace(tzinfo=datetime.now().astimezone().tzinfo)
    return parsed.astimezone(timezone.utc)


def _start_run_thread(run_id: str, req: "RunRequest", input_files: List[Path]) -> None:
    thread = threading.Thread(target=_execute_run, args=(run_id, req, input_files), daemon=True)
    thread.start()


def _safe_filename(name: str) -> str:
    cleaned = re.sub(r"[^A-Za-z0-9._-]+", "_", (name or "").strip())
    cleaned = cleaned.strip("._") or "file"
    return cleaned[:120]


def _ensure_dirs() -> None:
    for path in [DATA_DIR, UPLOAD_DIR, LOG_DIR, OUTPUT_DIR]:
        path.mkdir(parents=True, exist_ok=True)


def _read_json(path: Path, default):
    if not path.exists():
        return default
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        return default


def _write_json(path: Path, payload) -> None:
    tmp = path.with_suffix(path.suffix + ".tmp")
    tmp.write_text(json.dumps(payload, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")
    tmp.replace(path)


def _load_runs() -> Dict[str, dict]:
    return _read_json(RUNS_FILE, {})


def _save_runs(runs: Dict[str, dict]) -> None:
    _write_json(RUNS_FILE, runs)


def _get_run(run_id: str) -> Optional[dict]:
    with _STORE_LOCK:
        return _load_runs().get(run_id)


def _update_run(run_id: str, updates: dict) -> dict:
    with _STORE_LOCK:
        runs = _load_runs()
        record = runs.get(run_id)
        if record is None:
            raise KeyError(run_id)
        record.update(updates)
        runs[run_id] = record
        _save_runs(runs)
        return record


def _create_run(record: dict) -> dict:
    run_id = record["run_id"]
    with _STORE_LOCK:
        runs = _load_runs()
        runs[run_id] = record
        _save_runs(runs)
    return record


def _run_status_counts() -> Dict[str, int]:
    with _STORE_LOCK:
        runs = list(_load_runs().values())
    scheduled = sum(1 for r in runs if str(r.get("status") or "") == "scheduled")
    queued = sum(1 for r in runs if str(r.get("status") or "") == "queued")
    running = sum(1 for r in runs if str(r.get("status") or "") == "running")
    failed = sum(1 for r in runs if str(r.get("status") or "") == "failed")
    succeeded = sum(1 for r in runs if str(r.get("status") or "") in {"succeeded", "completed"})
    return {
        "total": len(runs),
        "scheduled": scheduled,
        "queued": queued,
        "running": running,
        "failed": failed,
        "succeeded": succeeded,
    }


def _load_uploads() -> Dict[str, dict]:
    return _read_json(UPLOADS_FILE, {})


def _save_uploads(uploads: Dict[str, dict]) -> None:
    _write_json(UPLOADS_FILE, uploads)


def _register_upload(record: dict) -> None:
    with _STORE_LOCK:
        uploads = _load_uploads()
        uploads[record["file_id"]] = record
        _save_uploads(uploads)


def _default_ui_config() -> dict:
    return {
        "standing_instruction": "",
        "knowledge_base_notes": "",
        "standing_instruction_attachments": [],
        "knowledge_base_attachments": [],
        "account_profile": {
            "operator_name": str(os.getenv("SF_OPERATOR_NAME", "Local Operator")).strip() or "Local Operator",
            "workspace": "SF Agentic AI",
        },
        "policy": {
            "credit_thresholds": dict(DEFAULT_CREDIT_THRESHOLDS),
            "credit_ratio_rules": [
                {
                    "id": "dscr_min",
                    "label": "DSCR min",
                    "description": "Coverage of debt service from cash flow",
                    "formula": "EBITDA / debt service",
                    "value": "1.25",
                    "locked": True,
                },
                {
                    "id": "icr_min",
                    "label": "ICR min",
                    "description": "Interest coverage ratio",
                    "formula": "EBIT / interest expense",
                    "value": "2.20",
                    "locked": True,
                },
                {
                    "id": "net_leverage_max",
                    "label": "Net leverage max",
                    "description": "Leverage after cash offset",
                    "formula": "(Debt - Cash) / EBITDA",
                    "value": "4.20",
                    "locked": True,
                },
                {
                    "id": "current_ratio_min",
                    "label": "Current ratio min",
                    "description": "Short-term liquidity coverage",
                    "formula": "Current assets / current liabilities",
                    "value": "1.10",
                    "locked": True,
                },
            ],
            "credit_decision_mode": "analysis_only",
            "report_length_words": 1800,
            "report_format": "credit_memo_v2",
            "sop_format": "standard_sop",
            "output_template_default": "autonomous",
            "condition_precedents": "\n".join(DEFAULT_CONDITION_PRECEDENTS),
            "adhoc_criteria": "",
            "structured_output_instruction": "",
            "market_intel_instruction": "",
            "market_intel_attachments": [],
        },
        "output_templates": dict(DEFAULT_OUTPUT_TEMPLATES),
        "rag": {
            "enabled": False,
            "instruction": "",
            "references": [],
            "attachments": [],
        },
        "google_drive": {
            "enabled": False,
            "source_folder_id": "",
            "output_folder": "/AgenticFlow/Outputs",
            "service_account_file": "",
        },
        "onedrive": {
            "enabled": False,
            "source_path": "",
            "output_path": "/SF-Agentic-AI/Outputs",
        },
        "telegram": {
            "enabled": bool(str(os.getenv("SF_NOTIFY_CHANNEL", "")).strip()),
            "channel": str(os.getenv("SF_NOTIFY_CHANNEL", "")).strip(),
        },
        "cloud_apis": _default_cloud_apis(),
        "defaults": {
            "audience": "credit_committee",
            "route_mode": "auto",
            "post_credit": "stop",
            "human_gate": "approve",
            "fast": True,
            "no_web_research": True,
            "autonomous": True,
            "output_formats": DEFAULT_OUTPUT_FORMATS,
        },
        "worker_registry": list(DEFAULT_WORKER_REGISTRY),
        "model_registry": list(DEFAULT_MODEL_REGISTRY),
    }


def _runtime_model_profile() -> dict:
    active_provider = _normalize_llm_provider(str(settings.llm_provider or "lm_studio"))
    label_suffix = "LM Studio" if active_provider == "lm_studio" else "Hermes Gateway"
    return {
        "id": str(settings.model or "local-model").strip() or "local-model",
        "label": f"{str(settings.model or 'local-model').strip()} ({label_suffix} active)",
        "provider": "local",
        "llm_provider": active_provider,
        "endpoint": str(settings.base_url or "").strip(),
        "enabled": True,
        "is_default": True,
    }


def _sanitize_ui_config(raw: dict) -> dict:
    base = _default_ui_config()
    if not isinstance(raw, dict):
        return base

    base["standing_instruction"] = str(raw.get("standing_instruction", "") or "").strip()
    base["knowledge_base_notes"] = str(raw.get("knowledge_base_notes", "") or "").strip()
    def _clean_attachment_list(value) -> List[str]:
        out: List[str] = []
        if isinstance(value, list):
            for item in value:
                if isinstance(item, dict):
                    fid = str(item.get("file_id", "") or "").strip()
                else:
                    fid = str(item or "").strip()
                if fid and fid not in out:
                    out.append(fid)
        return out
    base["standing_instruction_attachments"] = _clean_attachment_list(raw.get("standing_instruction_attachments", []))
    base["knowledge_base_attachments"] = _clean_attachment_list(raw.get("knowledge_base_attachments", []))

    account = raw.get("account_profile", {})
    if isinstance(account, dict):
        base["account_profile"]["operator_name"] = (
            str(account.get("operator_name", base["account_profile"]["operator_name"]) or "").strip()
            or "Local Operator"
        )
        base["account_profile"]["workspace"] = (
            str(account.get("workspace", base["account_profile"]["workspace"]) or "").strip()
            or "SF Agentic AI"
        )

    policy = raw.get("policy", {})
    if isinstance(policy, dict):
        thresholds = policy.get("credit_thresholds", {})
        if isinstance(thresholds, dict):
            cleaned_thresholds = dict(DEFAULT_CREDIT_THRESHOLDS)
            for key in DEFAULT_CREDIT_THRESHOLDS.keys():
                val = thresholds.get(key, cleaned_thresholds[key])
                try:
                    cleaned_thresholds[key] = float(val)
                except Exception:
                    pass
            base["policy"]["credit_thresholds"] = cleaned_thresholds
        ratio_rules = policy.get("credit_ratio_rules", [])
        if isinstance(ratio_rules, list):
            cleaned_rules = []
            for idx, row in enumerate(ratio_rules):
                if not isinstance(row, dict):
                    continue
                row_id = str(row.get("id", "") or "").strip()
                label = str(row.get("label", "") or "").strip()
                formula = str(row.get("formula", "") or "").strip()
                description = str(row.get("description", "") or "").strip()
                value = str(row.get("value", "") or "").strip()
                if not row_id or not label:
                    continue
                cleaned_rules.append({
                    "id": row_id,
                    "label": label,
                    "description": description,
                    "formula": formula,
                    "value": value,
                    "locked": bool(row.get("locked", False)),
                    "order": idx,
                })
            if cleaned_rules:
                base["policy"]["credit_ratio_rules"] = cleaned_rules
        try:
            base["policy"]["report_length_words"] = max(300, min(12000, int(policy.get("report_length_words", 1800))))
        except Exception:
            base["policy"]["report_length_words"] = 1800
        decision_mode = str(policy.get("credit_decision_mode", base["policy"]["credit_decision_mode"]) or "").strip()
        base["policy"]["credit_decision_mode"] = (
            decision_mode if decision_mode in ALLOWED_CREDIT_DECISION_MODES else "analysis_only"
        )
        base["policy"]["report_format"] = str(policy.get("report_format", base["policy"]["report_format"]) or "").strip() or "credit_memo_v2"
        base["policy"]["sop_format"] = str(policy.get("sop_format", base["policy"]["sop_format"]) or "").strip() or "standard_sop"
        output_template_default = str(policy.get("output_template_default", base["policy"]["output_template_default"]) or "").strip()
        base["policy"]["output_template_default"] = output_template_default or "autonomous"
        base["policy"]["condition_precedents"] = str(
            policy.get("condition_precedents", base["policy"]["condition_precedents"]) or ""
        ).strip()
        base["policy"]["adhoc_criteria"] = str(
            policy.get("adhoc_criteria", base["policy"]["adhoc_criteria"]) or ""
        ).strip()
        base["policy"]["structured_output_instruction"] = str(
            policy.get("structured_output_instruction", base["policy"]["structured_output_instruction"]) or ""
        ).strip()
        base["policy"]["market_intel_instruction"] = str(
            policy.get("market_intel_instruction", base["policy"]["market_intel_instruction"]) or ""
        ).strip()
        base["policy"]["market_intel_attachments"] = _clean_attachment_list(
            policy.get("market_intel_attachments", base["policy"].get("market_intel_attachments", []))
        )

    templates = raw.get("output_templates", {})
    if isinstance(templates, dict):
        cleaned_templates = dict(DEFAULT_OUTPUT_TEMPLATES)
        for key, default_text in DEFAULT_OUTPUT_TEMPLATES.items():
            cleaned_templates[key] = str(templates.get(key, default_text) or "").strip() or default_text
        for key, val in templates.items():
            k = str(key or "").strip()
            if not k:
                continue
            text = str(val or "").strip()
            if text:
                cleaned_templates[k] = text
        base["output_templates"] = cleaned_templates

    rag = raw.get("rag", {})
    if isinstance(rag, dict):
        base["rag"]["enabled"] = bool(rag.get("enabled", base["rag"]["enabled"]))
        base["rag"]["instruction"] = str(rag.get("instruction", base["rag"]["instruction"]) or "").strip()
        refs = rag.get("references", [])
        if isinstance(refs, list):
            cleaned_refs: List[str] = []
            for item in refs:
                val = str(item or "").strip()
                if val and val not in cleaned_refs:
                    cleaned_refs.append(val)
            base["rag"]["references"] = cleaned_refs[:80]
        base["rag"]["attachments"] = _clean_attachment_list(rag.get("attachments", base["rag"].get("attachments", [])))

    gd = raw.get("google_drive", {})
    if isinstance(gd, dict):
        base["google_drive"]["enabled"] = bool(gd.get("enabled", base["google_drive"]["enabled"]))
        base["google_drive"]["source_folder_id"] = str(gd.get("source_folder_id", "") or "").strip()
        base["google_drive"]["output_folder"] = str(gd.get("output_folder", base["google_drive"]["output_folder"]) or "").strip()
        base["google_drive"]["service_account_file"] = str(gd.get("service_account_file", "") or "").strip()

    od = raw.get("onedrive", {})
    if isinstance(od, dict):
        base["onedrive"]["enabled"] = bool(od.get("enabled", base["onedrive"]["enabled"]))
        base["onedrive"]["source_path"] = str(od.get("source_path", base["onedrive"]["source_path"]) or "").strip()
        base["onedrive"]["output_path"] = str(od.get("output_path", base["onedrive"]["output_path"]) or "").strip()

    tg = raw.get("telegram", {})
    if isinstance(tg, dict):
        base["telegram"]["enabled"] = bool(tg.get("enabled", base["telegram"]["enabled"]))
        base["telegram"]["channel"] = str(tg.get("channel", base["telegram"]["channel"]) or "").strip()

    cloud_apis = raw.get("cloud_apis", {})
    if isinstance(cloud_apis, dict):
        cleaned_cloud = _default_cloud_apis()
        for provider in SUPPORTED_CLOUD_API_PROVIDERS:
            entry = cloud_apis.get(provider, {})
            if not isinstance(entry, dict):
                continue
            cleaned_cloud[provider] = {
                "enabled": bool(entry.get("enabled", cleaned_cloud[provider]["enabled"])),
                "api_key": str(entry.get("api_key", "") or "").strip(),
                "endpoint": str(entry.get("endpoint", "") or "").strip(),
                "model": str(entry.get("model", "") or "").strip(),
            }
        base["cloud_apis"] = cleaned_cloud

    defaults = raw.get("defaults", {})
    if isinstance(defaults, dict):
        base["defaults"]["audience"] = str(defaults.get("audience", base["defaults"]["audience"]) or "").strip()
        base["defaults"]["route_mode"] = str(defaults.get("route_mode", base["defaults"]["route_mode"]) or "").strip()
        base["defaults"]["post_credit"] = str(defaults.get("post_credit", base["defaults"]["post_credit"]) or "").strip()
        base["defaults"]["human_gate"] = str(defaults.get("human_gate", base["defaults"]["human_gate"]) or "").strip()
        base["defaults"]["fast"] = bool(defaults.get("fast", base["defaults"]["fast"]))
        base["defaults"]["no_web_research"] = bool(defaults.get("no_web_research", base["defaults"]["no_web_research"]))
        base["defaults"]["autonomous"] = bool(defaults.get("autonomous", base["defaults"]["autonomous"]))
        formats = defaults.get("output_formats", base["defaults"]["output_formats"])
        if isinstance(formats, list):
            cleaned: List[str] = []
            for item in formats:
                fmt = str(item or "").strip().lower()
                if fmt in ALLOWED_OUTPUT_FORMATS and fmt not in cleaned:
                    cleaned.append(fmt)
            base["defaults"]["output_formats"] = cleaned or DEFAULT_OUTPUT_FORMATS

    workers = raw.get("worker_registry", [])
    if isinstance(workers, list):
        cleaned_workers: List[dict] = []
        seen_worker_ids: set[str] = set()
        for item in workers:
            if not isinstance(item, dict):
                continue
            wid = str(item.get("id", "") or "").strip()
            label = str(item.get("label", "") or "").strip()
            if not wid or wid in seen_worker_ids:
                continue
            seen_worker_ids.add(wid)
            cleaned_workers.append(
                {
                    "id": wid,
                    "label": label or wid,
                    "desc": str(item.get("desc", "") or "").strip(),
                    "enabled": bool(item.get("enabled", True)),
                    "category": str(item.get("category", "custom") or "custom").strip(),
                    "engage_when": str(item.get("engage_when", "") or "").strip(),
                    "default_intent": str(item.get("default_intent", "") or "").strip(),
                }
            )
        if cleaned_workers:
            for item in DEFAULT_WORKER_REGISTRY:
                wid = str(item.get("id", "") or "").strip()
                if not wid or wid in seen_worker_ids:
                    continue
                cleaned_workers.append(dict(item))
            base["worker_registry"] = cleaned_workers[:120]

    models = raw.get("model_registry", [])
    if isinstance(models, list):
        cleaned_models: List[dict] = []
        for item in models:
            if not isinstance(item, dict):
                continue
            mid = str(item.get("id", "") or "").strip()
            label = str(item.get("label", "") or "").strip()
            if not mid:
                continue
            temp_model = {
                "id": mid,
                "label": label or mid,
                "provider": str(item.get("provider", "local") or "local").strip(),
                "endpoint": str(item.get("endpoint", "") or "").strip(),
                "enabled": bool(item.get("enabled", True)),
                "is_default": bool(item.get("is_default", False)),
                "llm_provider": str(item.get("llm_provider", "") or "").strip(),
                "engage_when": str(item.get("engage_when", "") or "").strip(),
                "cost_hint_usd": str(item.get("cost_hint_usd", "") or "").strip(),
                "notes": str(item.get("notes", "") or "").strip(),
            }
            temp_model["llm_provider"] = _guess_llm_provider(temp_model)
            cleaned_models.append(temp_model)
        if cleaned_models:
            provider_match = any(
                str(m.get("llm_provider", "")).strip().lower() == str(settings.llm_provider or "").strip().lower()
                for m in cleaned_models
            )
            if not provider_match:
                runtime_model = _runtime_model_profile()
                for m in cleaned_models:
                    m["is_default"] = False
                cleaned_models.insert(0, runtime_model)
            default_idx = next((idx for idx, m in enumerate(cleaned_models) if m.get("is_default")), 0)
            for idx, model in enumerate(cleaned_models):
                model["is_default"] = idx == default_idx
            base["model_registry"] = cleaned_models[:120]

    if base["defaults"]["route_mode"] not in ALLOWED_ROUTE_MODES:
        base["defaults"]["route_mode"] = "auto"
    if base["defaults"]["post_credit"] not in ALLOWED_POST_CREDIT:
        base["defaults"]["post_credit"] = "stop"
    if base["defaults"]["human_gate"] not in ALLOWED_HUMAN_GATE:
        base["defaults"]["human_gate"] = "approve"
    return base


def _load_ui_config() -> dict:
    return _sanitize_ui_config(_read_json(UI_CONFIG_FILE, _default_ui_config()))


def _save_ui_config(payload: dict) -> dict:
    sanitized = _sanitize_ui_config(payload)
    _write_json(UI_CONFIG_FILE, sanitized)
    return sanitized


def _pick_default_model(model_registry: List[dict]) -> dict:
    if not isinstance(model_registry, list):
        return {}
    if not model_registry:
        return {}
    for item in model_registry:
        if isinstance(item, dict) and item.get("is_default"):
            return item
    return model_registry[0] if isinstance(model_registry[0], dict) else {}


def _bump_version_label(name: str, step: str = "rerun") -> str:
    raw = str(name or "").strip()
    if not raw:
        return "v1"
    raw = re.sub(r"(?:\s*\((?:rerun|copy|dup(?:licate)?)\))+\s*$", "", raw, flags=re.IGNORECASE).strip()
    m = re.search(r"(?:^|\s)v(\d+(?:\.\d+)*)\s*$", raw, re.IGNORECASE)
    if m:
        parts = [int(p) for p in m.group(1).split(".") if p.isdigit()]
        if parts:
            parts[-1] += 1
            return re.sub(r"(?:^|\s)v(\d+(?:\.\d+)*)\s*$", f" v{'.'.join(str(p) for p in parts)}", raw, flags=re.IGNORECASE).strip()
    return f"{raw} v1.1"


def _version_label_from_run_name(name: str, *, source_run_id: str = "") -> str:
    raw = str(name or "").strip()
    match = re.search(r"(?:^|\s)v(\d+(?:\.\d+)*)\s*$", raw, re.IGNORECASE)
    if match:
        return f"v{match.group(1)}"
    return "v1.1" if source_run_id else "v1"


def _guess_llm_provider(model: dict) -> str:
    if not isinstance(model, dict):
        return _normalize_llm_provider("")
    explicit = str(model.get("llm_provider", "") or "").strip()
    if explicit:
        return _normalize_llm_provider(explicit)
    provider_hint = str(model.get("provider", "") or "").strip().lower()
    endpoint = str(model.get("endpoint", "") or "").strip().lower()
    if "1234" in endpoint or endpoint.endswith("/v1"):
        return "lm_studio"
    if "18789" in endpoint or provider_hint in {"hermes", "hermes"}:
        return "hermes"
    if provider_hint == "cloud":
        return "hermes"
    if provider_hint == "local":
        return "lm_studio"
    return _normalize_llm_provider(provider_hint)


def _normalize_llm_provider(provider: str) -> str:
    raw = str(provider or "").strip().lower()
    if raw in {"hermes", "open_claw", "hermes"}:
        return "hermes"
    if raw in {"lm_studio", "lmstudio", "lm-studio"}:
        return "lm_studio"
    fallback = str(settings.llm_provider or "lm_studio").strip().lower()
    if fallback in {"hermes", "open_claw", "hermes"}:
        return "hermes"
    if fallback in {"lm_studio", "lmstudio", "lm-studio"}:
        return "lm_studio"
    return "lm_studio"


def _resolve_file_ids(file_ids: List[str]) -> List[Path]:
    with _STORE_LOCK:
        uploads = _load_uploads()
    paths: List[Path] = []
    missing: List[str] = []
    for file_id in file_ids:
        entry = uploads.get(file_id)
        if not entry:
            missing.append(file_id)
            continue
        path = Path(entry["stored_path"])
        if not path.exists():
            missing.append(file_id)
            continue
        paths.append(path)
    if missing:
        raise HTTPException(status_code=404, detail=f"Unknown or missing file_id(s): {', '.join(missing)}")
    return paths


def _attachments_to_text(file_ids: List[str], label: str) -> str:
    if not file_ids:
        return ""
    try:
        paths = _resolve_file_ids(file_ids)
    except HTTPException:
        return ""
    if not paths:
        return ""
    try:
        from tools.ocr import extract_text_from_files
        result = extract_text_from_files([str(p) for p in paths])
        text = str(result.get("text", "") or "").strip()
        if text:
            return f"[{label} attachments]\n" + text
    except Exception:
        pass
    names = ", ".join(p.name for p in paths[:20])
    return f"[{label} attachments]\n{names}"


def _build_run_cmd(req: "RunRequest", input_files: List[Path]) -> List[str]:
    output_family = _output_family_for_request(req)
    cmd = [
        sys.executable,
        "main.py",
        "--intent",
        req.intent,
        "--audience",
        req.audience,
        "--route-mode",
        req.route_mode,
        "--post-credit",
        req.post_credit,
        "--human-gate",
        req.human_gate,
        "--report-format",
        output_family,
        "--sop-format",
        req.sop_format,
        "--output-template",
        req.output_template_id or "autonomous",
        "--structured-output",
        req.structured_output_instruction,
        "--report-length-words",
        str(req.report_length_words or 0),
        "--credit-decision-mode",
        req.credit_decision_mode,
    ]

    if req.workers:
        cmd.extend(["--workers", ",".join(req.workers)])

    notes: List[str] = []
    if req.standing_instruction.strip():
        notes.append(f"[Standing Instruction]\n{req.standing_instruction.strip()}")
    standing_attachment_text = _attachments_to_text(getattr(req, "standing_instruction_attachment_ids", []) or [], "Standing instruction")
    if standing_attachment_text:
        notes.append(standing_attachment_text)
    if req.knowledge_base_notes.strip():
        notes.append(f"[Knowledge Base]\n{req.knowledge_base_notes.strip()}")
    knowledge_attachment_text = _attachments_to_text(getattr(req, "knowledge_base_attachment_ids", []) or [], "Knowledge base")
    if knowledge_attachment_text:
        notes.append(knowledge_attachment_text)
    if req.manager_instruction.strip():
        notes.append(f"[Run-specific Instruction]\n{req.manager_instruction.strip()}")
    if req.market_intel_instruction.strip():
        notes.append(f"[Market Intel Instruction]\n{req.market_intel_instruction.strip()}")
    if req.report_length_words > 0:
        notes.append(f"[Report Length Target]\n{req.report_length_words} words")
    if req.report_format.strip():
        notes.append(f"[Report Format]\n{req.report_format.strip()}")
    if req.sop_format.strip():
        notes.append(f"[SOP Format]\n{req.sop_format.strip()}")
    if output_family == "credit_memo":
        notes.append(f"[Credit Memo Objective]\n{req.credit_decision_mode}")
    else:
        notes.append(f"[Output Family]\n{output_family}")
    if req.output_template_id.strip():
        notes.append(f"[Output Template]\n{req.output_template_id.strip()}")
    if req.output_template_id.strip() == "custom":
        notes.append("[Structured Template Mode]\ncustom")
    if req.drive_source_link.strip():
        notes.append(f"[Google Drive Source Link]\n{req.drive_source_link.strip()}")
    if req.onedrive_source_link.strip():
        notes.append(f"[OneDrive Source Link]\n{req.onedrive_source_link.strip()}")
    custom_intent = req.custom_intent.strip()
    if req.intent == "custom" and custom_intent:
        notes.append(f"[Custom Intent]\n{custom_intent}")
    if req.credit_decision_mode == "decision_recommendation" and req.condition_precedents.strip():
        notes.append(f"[Condition Precedents]\n{req.condition_precedents.strip()}")
    if req.adhoc_criteria.strip():
        notes.append(f"[Ad-hoc Criteria]\n{req.adhoc_criteria.strip()}")
    if req.structured_output_instruction.strip():
        notes.append(f"[Structured Output Requirement]\n{req.structured_output_instruction.strip()}")
    if req.rag_instruction.strip():
        notes.append(f"[RAG / Knowledge References]\n{req.rag_instruction.strip()}")
    if req.knowledge_file_ids:
        notes.append(f"[Knowledge Attachments]\n{len(req.knowledge_file_ids)} file(s) attached as reference pack")
    if req.credit_decision_mode == "decision_recommendation" and req.credit_thresholds:
        ratio_lines: List[str] = []
        for key, val in req.credit_thresholds.items():
            ratio_lines.append(f"- {key}: {val}")
        if ratio_lines:
            notes.append("[Credit Ratio Thresholds]\n" + "\n".join(ratio_lines))
    if req.extra_thresholds:
        extra_lines: List[str] = []
        for row in req.extra_thresholds:
            if not isinstance(row, dict):
                continue
            label = str(row.get("label") or row.get("id") or "extra").strip()
            value = str(row.get("value") or "").strip()
            formula = str(row.get("formula") or "").strip()
            extra_lines.append(f"- {label}: {value}" + (f" ({formula})" if formula else ""))
        if extra_lines:
            notes.append("[Extra Credit Thresholds]\n" + "\n".join(extra_lines))
    if notes:
        cmd.extend(["--manager-note", "\n\n".join(notes)])
    if input_files:
        cmd.extend(["--files", *[str(p) for p in input_files]])
    elif req.text.strip():
        cmd.extend(["--text", req.text.strip()])
    else:
        cmd.append("--sample")

    if req.autonomous:
        cmd.append("--autonomous")
    if req.fast:
        cmd.append("--fast")
    if req.skip_preflight:
        cmd.append("--skip-preflight")
    if req.no_web_research:
        cmd.append("--no-web-research")
    return cmd


def _snapshot_outputs() -> Dict[str, float]:
    snapshot: Dict[str, float] = {}
    if not OUTPUT_DIR.exists():
        return snapshot
    for p in OUTPUT_DIR.iterdir():
        if p.is_file():
            snapshot[str(p.resolve())] = p.stat().st_mtime
    return snapshot


def _collect_artifacts(before: Dict[str, float]) -> List[dict]:
    artifacts: List[dict] = []
    if not OUTPUT_DIR.exists():
        return artifacts
    for p in sorted(OUTPUT_DIR.iterdir(), key=lambda item: item.stat().st_mtime, reverse=True):
        if not p.is_file():
            continue
        key = str(p.resolve())
        old_mtime = before.get(key)
        if old_mtime is not None and old_mtime >= p.stat().st_mtime:
            continue
        artifacts.append(
            {
                "name": p.name,
                "path": key,
                "size_bytes": p.stat().st_size,
                "mtime": datetime.fromtimestamp(p.stat().st_mtime, timezone.utc).isoformat(),
                "mime_type": _detect_mime_type(p),
            }
        )
    return artifacts


def _detect_mime_type(path: Path) -> str:
    guess, _ = mimetypes.guess_type(path.name)
    return guess or "application/octet-stream"


def _file_sha256(path: Path) -> str:
    if not path.exists() or not path.is_file():
        return ""
    digest = hashlib.sha256()
    try:
        with path.open("rb") as handle:
            for chunk in iter(lambda: handle.read(1024 * 1024), b""):
                digest.update(chunk)
        return digest.hexdigest()
    except Exception:
        return ""


def _artifact_id(run_id: str, name: str) -> str:
    stable = uuid.uuid5(uuid.NAMESPACE_URL, f"sf-agentic-ai/artifact/{run_id}/{name}")
    return f"art_{stable.hex[:16]}"


def _artifact_version_label(record: Optional[dict]) -> str:
    if not isinstance(record, dict):
        return "v1"
    explicit = str(record.get("version_label") or "").strip()
    if explicit:
        return explicit
    request = record.get("request") if isinstance(record.get("request"), dict) else {}
    run_name = str(request.get("run_name") or "").strip()
    return _version_label_from_run_name(run_name, source_run_id=str(record.get("source_run_id") or ""))


def _artifact_kind(item: dict, path: Path, record: Optional[dict] = None) -> str:
    name = str(item.get("name") or path.name or "").lower()
    generated_by = str(item.get("generated_by") or "").lower()
    if name.endswith(".metadata.json") or name.endswith("_metadata.json"):
        return "metadata"
    if "postcredit" in name or "post_credit" in name:
        return "post_credit"
    if name.endswith(".log") or "warning" in name or generated_by == "export:warnings":
        return "log"

    # Intent-aware classification for ambiguous prefixes
    intent = ""
    if record and isinstance(record.get("request"), dict):
        intent = str(record["request"].get("intent") or "").lower()
    
    family_from_intent = OUTPUT_FAMILY_KINDS.get(intent)

    if name.startswith("cc_memo") or name.startswith("creditmemo"):
        return family_from_intent or "credit_memo"
    if name.startswith("analysisbrief"):
        return family_from_intent or "analysis_brief"
    if name.startswith("sop"):
        return family_from_intent or "sop"
    if name.startswith("meetingminutes"):
        return family_from_intent or "meeting_minutes"
    if name.startswith("pipelinereport"):
        return family_from_intent or "pipeline_report"
    if name.startswith("dataanalysis"):
        return family_from_intent or "data_analysis"
    if name.startswith("marketintel"):
        return family_from_intent or "market_intel"
    
    if generated_by.startswith("export:"):
        return str(item.get("kind") or "other")
    return "other"


def _artifact_stable_name(run_id: str, record: Optional[dict], item: dict, path: Path, index: int) -> str:
    existing = str(item.get("stable_name") or "").strip()
    if existing:
        return existing
    request = record.get("request") if isinstance(record, dict) and isinstance(record.get("request"), dict) else {}
    prefix = str(request.get("run_name") or path.stem or item.get("name") or "artifact").strip()
    prefix = re.sub(r"(?:^|\s)v\d+(?:\.\d+)*\s*$", "", prefix, flags=re.IGNORECASE).strip() or prefix
    version = _artifact_version_label(record).replace(".", "_")
    ext = path.suffix or Path(str(item.get("name") or "")).suffix
    return f"{_safe_filename(prefix)}_{version}_{run_id[:8]}_{index + 1:02d}{ext.lower()}"


def _decorate_artifact(run_id: str, item: dict, record: Optional[dict] = None, index: int = 0) -> dict:
    raw_name = str(item.get("name") or "").strip()
    raw_path = str(item.get("path") or "").strip()
    path = Path(raw_path) if raw_path else Path(raw_name)
    name = raw_name or path.name
    if not name:
        name = f"artifact_{index + 1}"
    if not raw_path:
        path = OUTPUT_DIR / name

    exists = path.exists() and path.is_file()
    stat = path.stat() if exists else None
    mtime = (
        datetime.fromtimestamp(stat.st_mtime, timezone.utc).isoformat()
        if stat
        else str(item.get("mtime") or item.get("created_at") or "")
    )
    size_bytes = int(stat.st_size) if stat else int(item.get("size_bytes") or 0)
    path_status = "available" if exists else "missing"
    validation_status = "ok"
    if exists and size_bytes == 0:
        path_status = "empty"
        validation_status = "invalid_empty"
    elif not exists:
        validation_status = "missing"
    created_at = str(item.get("created_at") or mtime or (record or {}).get("finished_at") or (record or {}).get("created_at") or "")
    mime_type = str(item.get("mime_type") or _detect_mime_type(path))
    version_label = str(item.get("version_label") or _artifact_version_label(record))
    reported_path_status = str(item.get("path_status") or path_status)
    if validation_status == "invalid_empty":
        reported_path_status = "empty"
    decorated = {
        **item,
        "artifact_schema_version": int(item.get("artifact_schema_version") or ARTIFACT_SCHEMA_VERSION),
        "artifact_id": str(item.get("artifact_id") or _artifact_id(run_id, name)),
        "run_id": run_id,
        "name": name,
        "display_name": str(item.get("display_name") or name),
        "stable_name": _artifact_stable_name(run_id, record, item, path, index),
        "kind": str(item.get("kind") or _artifact_kind(item, path, record=record)),
        "version": version_label,
        "version_label": version_label,
        "source_run_id": str((record or {}).get("source_run_id") or item.get("source_run_id") or ""),
        "path": str(path.resolve()) if exists else raw_path,
        "path_status": reported_path_status,
        "validation_status": str(item.get("validation_status") or validation_status),
        "size_bytes": size_bytes,
        "mtime": mtime,
        "created_at": created_at,
        "mime_type": mime_type,
        "sha256": str(item.get("sha256") or (_file_sha256(path) if exists else "")),
        "download_url": f"/api/artifacts/{run_id}/{str(item.get('artifact_id') or _artifact_id(run_id, name))}",
        "preview_url": f"/api/artifacts/{run_id}/{str(item.get('artifact_id') or _artifact_id(run_id, name))}/preview",
    }
    return decorated


def _decorate_artifacts(run_id: str, artifacts: List[dict], record: Optional[dict] = None) -> List[dict]:
    decorated: List[dict] = []
    for idx, item in enumerate(artifacts or []):
        if isinstance(item, dict):
            decorated.append(_decorate_artifact(run_id, item, record=record, index=idx))
    return decorated


def _artifact_priority(item: dict) -> tuple:
    kind = str(item.get("kind") or "").lower()
    validation = str(item.get("validation_status") or "").lower()
    path_status = str(item.get("path_status") or "").lower()
    size = int(item.get("size_bytes") or 0)
    created = str(item.get("created_at") or item.get("mtime") or "")
    invalid = validation not in {"", "ok"} or path_status in {"empty", "missing"} or size <= 0
    kind_rank = {
        "credit_memo": 0,
        "analysis_brief": 0,
        "sop": 0,
        "meeting_minutes": 0,
        "pipeline_report": 0,
        "data_analysis": 0,
        "market_intel": 0,
        "post_credit": 1,
        "other": 2,
        "log": 3,
        "metadata": 4,
    }.get(kind, 2)
    return (1 if invalid else 0, kind_rank, -size, created)


def _decorate_run(record: dict) -> dict:
    run_id = str(record.get("run_id") or "")
    decorated = dict(record)
    decorated["version_label"] = _artifact_version_label(record)
    decorated["lineage_root_run_id"] = str(record.get("lineage_root_run_id") or record.get("source_run_id") or run_id)
    artifacts = _decorate_artifacts(run_id, list(record.get("artifacts") or []), record=record)
    decorated["artifacts"] = artifacts
    decorated["artifacts_count"] = len(artifacts)
    decorated["artifacts_url"] = f"/api/runs/{run_id}/artifacts" if run_id else ""
    decorated["status_url"] = f"/api/runs/{run_id}" if run_id else ""
    decorated["latest_artifact"] = sorted(artifacts, key=_artifact_priority)[0] if artifacts else None
    return decorated


def _persist_decorated_artifacts(record: dict) -> List[dict]:
    run_id = str(record.get("run_id") or "")
    artifacts = _decorate_artifacts(run_id, list(record.get("artifacts") or []), record=record)
    if artifacts != record.get("artifacts"):
        _update_run(run_id, {"artifacts": artifacts, "artifacts_count": len(artifacts)})
    return artifacts


def _find_artifact(record: dict, identifier: str) -> Optional[dict]:
    run_id = str(record.get("run_id") or "")
    target = str(identifier or "").strip()
    for item in _decorate_artifacts(run_id, list(record.get("artifacts") or []), record=record):
        matches = {
            str(item.get("artifact_id") or ""),
            str(item.get("name") or ""),
            str(item.get("stable_name") or ""),
            str(item.get("display_name") or ""),
        }
        if target in matches:
            return item
    return None


def _normalize_output_formats(formats: List[str]) -> List[str]:
    cleaned: List[str] = []
    for item in formats:
        fmt = str(item or "").strip().lower()
        if fmt in ALLOWED_OUTPUT_FORMATS and fmt not in cleaned:
            cleaned.append(fmt)
    return cleaned


def _pick_primary_text_artifact(artifacts: List[dict]) -> Optional[Path]:
    preferred = [".md", ".txt", ".json", ".csv", ".tsv"]
    by_ext: Dict[str, List[Tuple[int, Path]]] = {ext: [] for ext in preferred}
    for item in artifacts:
        raw = str(item.get("path", "")).strip()
        if not raw:
            continue
        p = Path(raw)
        if not p.exists() or not p.is_file():
            continue
        ext = p.suffix.lower()
        if ext in by_ext:
            by_ext[ext].append((int(item.get("size_bytes", 0) or 0), p))
    for ext in preferred:
        if by_ext[ext]:
            by_ext[ext].sort(key=lambda x: x[0], reverse=True)
            return by_ext[ext][0][1]
    return None


def _read_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def _normalize_text_for_export(raw: str) -> str:
    text = str(raw or "").replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _unique_output_path(base_stem: str, ext: str) -> Path:
    safe_stem = _safe_filename(base_stem) or "output"
    safe_ext = ext.lstrip(".").lower()
    candidate = OUTPUT_DIR / f"{safe_stem}.{safe_ext}"
    if not candidate.exists():
        return candidate
    suffix = datetime.now().strftime("%Y%m%d_%H%M%S_%f")
    return OUTPUT_DIR / f"{safe_stem}_{suffix}.{safe_ext}"


def _markdown_sections(text: str) -> List[Tuple[str, List[str]]]:
    sections: List[Tuple[str, List[str]]] = []
    current_title = "Content"
    current_lines: List[str] = []
    for line in text.splitlines():
        stripped = line.strip()
        if stripped.startswith("#"):
            if current_lines:
                sections.append((current_title, current_lines))
            current_title = stripped.lstrip("#").strip() or "Content"
            current_lines = []
            continue
        current_lines.append(line.rstrip())
    if current_lines:
        sections.append((current_title, current_lines))
    if not sections:
        sections = [("Content", [text])]
    return sections


def _split_multiline_items(raw: str) -> List[str]:
    items: List[str] = []
    for line in str(raw or "").splitlines():
        stripped = line.strip().lstrip("-").strip()
        if stripped and stripped not in items:
            items.append(stripped)
    return items


def _template_instruction(req: "RunRequest") -> str:
    parts: List[str] = []
    if req.output_template_id.strip():
        if req.output_template_id.strip() in {"autonomous", "not_applicable"}:
            parts.append(f"Template: {req.output_template_id.strip()} (no enforced structure)")
        else:
            parts.append(f"Template: {req.output_template_id.strip()}")
    if req.report_format.strip():
        parts.append(f"Report format: {req.report_format.strip()}")
    if req.sop_format.strip():
        parts.append(f"SOP format: {req.sop_format.strip()}")
    if _output_family_for_request(req) == "credit_memo" and req.credit_decision_mode.strip():
        parts.append(f"Credit objective: {req.credit_decision_mode.strip()}")
    if req.structured_output_instruction.strip():
        parts.append(req.structured_output_instruction.strip())
    return " | ".join(parts)


def _write_docx(path: Path, title: str, text: str, req: Optional["RunRequest"] = None) -> None:
    if Document is None or Pt is None:
        raise RuntimeError("python-docx is not available")
    doc = Document()
    normal = doc.styles["Normal"]
    normal.font.name = "Calibri"
    normal.font.size = Pt(11)
    doc.add_heading(title, level=0)
    profile_line = _template_instruction(req) if req else ""
    if profile_line:
        p = doc.add_paragraph(profile_line)
        p.runs[0].italic = True

    for line in text.splitlines():
        stripped = line.strip()
        if not stripped:
            doc.add_paragraph("")
            continue
        if stripped.startswith("#"):
            level = min(3, max(1, stripped.count("#")))
            doc.add_heading(stripped.lstrip("#").strip(), level=level)
            continue
        if stripped.startswith("- ") or stripped.startswith("* "):
            doc.add_paragraph(stripped[2:].strip(), style="List Bullet")
            continue
        doc.add_paragraph(line.rstrip())

    if req:
        is_credit = _output_family_for_request(req) == "credit_memo"
        cp_items = _split_multiline_items(req.condition_precedents)
        if is_credit and req.credit_decision_mode == "decision_recommendation" and cp_items:
            doc.add_heading("Condition Precedents", level=1)
            for item in cp_items:
                doc.add_paragraph(item, style="List Bullet")
        if req.adhoc_criteria.strip():
            doc.add_heading("Custom Criteria", level=1)
            doc.add_paragraph(req.adhoc_criteria.strip())
    doc.save(str(path))


def _write_xlsx(path: Path, title: str, text: str, req: Optional["RunRequest"] = None) -> None:
    if Workbook is None:
        raise RuntimeError("openpyxl is not available")
    wb = Workbook()
    ws = wb.active
    ws.title = "Summary"
    ws["A1"] = "Section"
    ws["B1"] = "Content"
    if Font is not None:
        ws["A1"].font = Font(bold=True)
        ws["B1"].font = Font(bold=True)
    ws.column_dimensions["A"].width = 36
    ws.column_dimensions["B"].width = 120
    row = 2
    for section, lines in _markdown_sections(text):
        content = "\n".join([ln for ln in lines if ln.strip()]).strip()
        if not content:
            continue
        ws.cell(row=row, column=1, value=section)
        ws.cell(row=row, column=2, value=content)
        if Alignment is not None:
            ws.cell(row=row, column=2).alignment = Alignment(wrap_text=True, vertical="top")
        row += 1

    raw = wb.create_sheet("Raw Text")
    raw["A1"] = title
    if Font is not None:
        raw["A1"].font = Font(bold=True)
    raw["A2"] = text
    if Alignment is not None:
        raw["A2"].alignment = Alignment(wrap_text=True, vertical="top")
    raw.column_dimensions["A"].width = 140

    if req:
        is_credit = _output_family_for_request(req) == "credit_memo"
        meta = wb.create_sheet("Decision Controls" if is_credit else "Output Controls")
        meta["A1"] = "Field"
        meta["B1"] = "Value"
        if Font is not None:
            meta["A1"].font = Font(bold=True)
            meta["B1"].font = Font(bold=True)
        rows = [
            ("Output template", req.output_template_id),
            ("Report format", req.report_format),
            ("SOP format", req.sop_format),
            ("Credit objective", req.credit_decision_mode if is_credit else "not_applicable"),
            ("Report length words", req.report_length_words),
            ("Ad-hoc criteria", req.adhoc_criteria),
        ]
        r = 2
        for key, val in rows:
            meta.cell(row=r, column=1, value=key)
            meta.cell(row=r, column=2, value=str(val or ""))
            r += 1
        ratios = req.credit_thresholds if isinstance(req.credit_thresholds, dict) else {}
        if is_credit and req.credit_decision_mode == "decision_recommendation" and ratios:
            meta.cell(row=r, column=1, value="Credit thresholds")
            r += 1
            for k, v in ratios.items():
                meta.cell(row=r, column=1, value=str(k))
                meta.cell(row=r, column=2, value=str(v))
                r += 1
        cps = _split_multiline_items(req.condition_precedents)
        if is_credit and req.credit_decision_mode == "decision_recommendation" and cps:
            meta.cell(row=r, column=1, value="Condition precedents")
            r += 1
            for cp in cps:
                meta.cell(row=r, column=2, value=cp)
                r += 1
        meta.column_dimensions["A"].width = 38
        meta.column_dimensions["B"].width = 120
    wb.save(str(path))


def _write_pptx(path: Path, title: str, text: str) -> None:
    if Presentation is None or Inches is None or PptPt is None:
        raise RuntimeError("python-pptx is not available")
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = f"Generated {datetime.now().strftime('%Y-%m-%d %H:%M')}"

    sections = _markdown_sections(text)
    for section, lines in sections[:12]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = section[:120]
        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()
        bullet_lines: List[str] = []
        for raw_line in lines:
            stripped = raw_line.strip()
            if not stripped:
                continue
            if stripped.startswith("- ") or stripped.startswith("* "):
                bullet_lines.append(stripped[2:].strip())
            else:
                bullet_lines.extend(textwrap.wrap(stripped, width=90) or [stripped])
            if len(bullet_lines) >= 10:
                break
        if not bullet_lines:
            bullet_lines = ["No additional detail."]
        for idx, line in enumerate(bullet_lines[:10]):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = line
            p.level = 0
            p.font.size = PptPt(18 if idx == 0 else 16)
    prs.save(str(path))


def _pdf_escape(line: str) -> str:
    return line.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")


def _write_minimal_pdf(path: Path, title: str, text: str) -> None:
    width = 612
    height = 792
    margin_x = 54
    start_y = height - 54
    line_height = 13.5
    max_chars = 94
    max_lines_per_page = 48

    lines: List[str] = [title, ""]
    for raw_line in text.splitlines():
        if not raw_line.strip():
            lines.append("")
            continue
        wrapped = textwrap.wrap(raw_line, width=max_chars) or [raw_line]
        lines.extend(wrapped)

    pages: List[List[str]] = []
    cursor: List[str] = []
    for line in lines:
        cursor.append(line)
        if len(cursor) >= max_lines_per_page:
            pages.append(cursor)
            cursor = []
    if cursor:
        pages.append(cursor)
    if not pages:
        pages = [[title, "(empty)"]]

    objects: List[bytes] = []
    objects.append(b"<< /Type /Catalog /Pages 2 0 R >>")
    page_kids_refs = []

    page_objs: List[bytes] = []
    content_objs: List[bytes] = []

    base_obj_num = 3
    for i, page_lines in enumerate(pages):
        page_obj_num = base_obj_num + (i * 2)
        content_obj_num = page_obj_num + 1
        page_kids_refs.append(f"{page_obj_num} 0 R")
        stream_lines = [
            "BT",
            "/F1 11 Tf",
            f"{line_height:.2f} TL",
            f"{margin_x} {start_y} Td",
        ]
        first = True
        for line in page_lines:
            escaped = _pdf_escape(line)
            if first:
                stream_lines.append(f"({escaped}) Tj")
                first = False
            else:
                stream_lines.append("T*")
                stream_lines.append(f"({escaped}) Tj")
        stream_lines.append("ET")
        content = "\n".join(stream_lines).encode("latin-1", errors="replace")
        content_obj = (
            f"<< /Length {len(content)} >>\nstream\n".encode("ascii")
            + content
            + b"\nendstream"
        )
        page_obj = (
            f"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 {width} {height}] "
            f"/Resources << /Font << /F1 3 0 R >> >> /Contents {content_obj_num} 0 R >>"
        ).encode("ascii")
        page_objs.append(page_obj)
        content_objs.append(content_obj)

    objects.append(f"<< /Type /Pages /Kids [{' '.join(page_kids_refs)}] /Count {len(page_kids_refs)} >>".encode("ascii"))
    objects.append(b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>")
    for i in range(len(page_objs)):
        objects.append(page_objs[i])
        objects.append(content_objs[i])

    output = bytearray()
    output.extend(b"%PDF-1.4\n")
    xref_positions = [0]
    for idx, obj in enumerate(objects, start=1):
        xref_positions.append(len(output))
        output.extend(f"{idx} 0 obj\n".encode("ascii"))
        output.extend(obj)
        output.extend(b"\nendobj\n")
    xref_start = len(output)
    output.extend(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    output.extend(b"0000000000 65535 f \n")
    for pos in xref_positions[1:]:
        output.extend(f"{pos:010d} 00000 n \n".encode("ascii"))
    output.extend(
        (
            f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
            f"startxref\n{xref_start}\n%%EOF\n"
        ).encode("ascii")
    )
    path.write_bytes(bytes(output))


def _build_generated_artifact(path: Path, generated_by: str, source_path: Optional[Path] = None, kind: str = "") -> dict:
    return {
        "name": path.name,
        "path": str(path.resolve()),
        "size_bytes": path.stat().st_size,
        "mtime": datetime.fromtimestamp(path.stat().st_mtime, timezone.utc).isoformat(),
        "generated": True,
        "generated_by": generated_by,
        "source_path": str(source_path.resolve()) if source_path else "",
        "mime_type": _detect_mime_type(path),
        "kind": kind or "",
    }


def _generate_requested_artifacts(run_id: str, req: "RunRequest", artifacts: List[dict]) -> List[dict]:
    requested = _normalize_output_formats(req.output_formats or [])
    if not requested:
        requested = DEFAULT_OUTPUT_FORMATS

    primary = _pick_primary_text_artifact(artifacts)
    if not primary:
        return []

    text = _normalize_text_for_export(_read_text(primary))
    if not text:
        return []

    base_stem = Path(primary.stem).name
    run_suffix = run_id[:8]
    output_family = _output_family_for_request(req)
    artifact_kind = OUTPUT_FAMILY_KINDS.get(output_family, "other")
    version_slug = _safe_filename(_version_label_from_run_name(req.run_name)).replace(".", "_") or "v1"
    title = req.run_name.strip() or f"{output_family.replace('_', ' ').title()} {run_suffix}"

    generated: List[dict] = []
    skips: List[str] = []
    existing_names = {str(item.get("name", "")).strip() for item in artifacts}

    def make_target(ext: str) -> Optional[Path]:
        candidate = _unique_output_path(f"{base_stem}_{version_slug}_{run_suffix}", ext)
        if candidate.name in existing_names:
            candidate = _unique_output_path(f"{base_stem}_{version_slug}_{run_suffix}_rev", ext)
        return candidate

    for fmt in requested:
        ext = fmt.lower()
        if ext == primary.suffix.lstrip(".").lower():
            continue
        try:
            target = make_target(ext)
            if target is None:
                continue
            if ext == "txt":
                target.write_text(text + "\n", encoding="utf-8")
            elif ext == "md":
                target.write_text(text + "\n", encoding="utf-8")
            elif ext == "docx":
                _write_docx(target, title, text, req=req)
            elif ext == "xlsx":
                _write_xlsx(target, title, text, req=req)
            elif ext == "pptx":
                _write_pptx(target, title, text)
            elif ext == "pdf":
                _write_minimal_pdf(target, title, text)
            else:
                continue
            generated.append(_build_generated_artifact(target, generated_by=f"export:{ext}", source_path=primary, kind=artifact_kind))
        except Exception as exc:
            skips.append(f"{fmt}: {exc}")

    if skips:
        warn_path = _unique_output_path(f"{base_stem}_{run_suffix}_export_warnings", "txt")
        warn_path.write_text(
            "Some requested formats were skipped:\n\n" + "\n".join(f"- {msg}" for msg in skips) + "\n",
            encoding="utf-8",
        )
        generated.append(_build_generated_artifact(warn_path, generated_by="export:warnings", source_path=primary, kind="log"))

    return generated


def _artifact_preview(path: Path, max_chars: int = 12000) -> dict:
    ext = path.suffix.lower()
    if ext in {".txt", ".md", ".json", ".csv", ".tsv", ".log"}:
        content = path.read_text(encoding="utf-8", errors="ignore")
        truncated = len(content) > max_chars
        return {
            "kind": "text",
            "content": content[:max_chars],
            "truncated": truncated,
            "mime_type": _detect_mime_type(path),
        }
    if ext == ".docx":
        if Document is None:
            raise RuntimeError("python-docx is not available")
        doc = Document(str(path))
        content = "\n".join([p.text for p in doc.paragraphs if p.text is not None])
        truncated = len(content) > max_chars
        return {
            "kind": "text",
            "content": content[:max_chars],
            "truncated": truncated,
            "mime_type": _detect_mime_type(path),
        }
    if ext == ".xlsx":
        if load_workbook is None:
            raise RuntimeError("openpyxl is not available")
        wb = load_workbook(filename=str(path), read_only=True, data_only=True)
        ws = wb[wb.sheetnames[0]]
        rows: List[str] = []
        for idx, row in enumerate(ws.iter_rows(values_only=True)):
            if idx >= 120:
                break
            row_values = ["" if v is None else str(v) for v in row]
            if any(val.strip() for val in row_values):
                rows.append(" | ".join(row_values[:12]))
        wb.close()
        content = "\n".join(rows) if rows else "(No visible cells in first sheet)"
        truncated = len(rows) >= 120
        return {
            "kind": "table",
            "content": content[:max_chars],
            "truncated": truncated,
            "mime_type": _detect_mime_type(path),
        }
    if ext == ".pptx":
        if Presentation is None:
            return {
                "kind": "unsupported",
                "content": "Preview unavailable: install python-pptx to parse slide text.",
                "truncated": False,
                "mime_type": _detect_mime_type(path),
            }
        prs = Presentation(str(path))
        chunks: List[str] = []
        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_lines: List[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_lines.extend([ln.strip() for ln in shape.text.splitlines() if ln.strip()])
            if slide_lines:
                chunks.append(f"[Slide {slide_idx}] " + " | ".join(slide_lines[:6]))
        content = "\n".join(chunks) if chunks else "(No text found in slides)"
        truncated = len(content) > max_chars
        return {
            "kind": "text",
            "content": content[:max_chars],
            "truncated": truncated,
            "mime_type": _detect_mime_type(path),
        }
    return {
        "kind": "unsupported",
        "content": f"Preview unavailable for {ext or 'this'} file type. Use download instead.",
        "truncated": False,
        "mime_type": _detect_mime_type(path),
    }


def _execute_run(run_id: str, req: "RunRequest", input_files: List[Path]) -> None:
    stdout_path = LOG_DIR / f"{run_id}.stdout.log"
    stderr_path = LOG_DIR / f"{run_id}.stderr.log"
    before_outputs = _snapshot_outputs()
    cmd = _build_run_cmd(req, input_files)

    _update_run(
        run_id,
        {
            "status": "running",
            "started_at": _now_iso(),
            "cmd": cmd,
            "stdout_log": str(stdout_path.resolve()),
            "stderr_log": str(stderr_path.resolve()),
        },
    )

    env = dict(os.environ)
    env.setdefault("PYTHONPYCACHEPREFIX", "/tmp/pycache")
    # Suppress noisy LibreSSL urllib3 warning in subprocess logs on macOS Python builds.
    env["PYTHONWARNINGS"] = "ignore:urllib3 v2 only supports OpenSSL 1.1.1+"
    provider = _normalize_llm_provider(req.llm_provider)
    model = str(req.llm_model or "").strip()
    base_url = str(req.llm_base_url or "").strip()
    if provider:
        env["LLM_PROVIDER"] = provider
    if model:
        if provider == "hermes":
            env["HERMES_MODEL"] = model
            env["HERMES_MODEL"] = model
        else:
            env["LM_STUDIO_MODEL"] = model
    if base_url:
        if provider == "hermes":
            env["HERMES_BASE_URL"] = base_url
            env["HERMES_BASE_URL"] = base_url
        else:
            env["LM_STUDIO_BASE_URL"] = base_url
    if int(req.timeout_sec or 0) > 0:
        env["LLM_TIMEOUT_SEC"] = str(int(req.timeout_sec))

    with stdout_path.open("w", encoding="utf-8") as stdout_file, stderr_path.open("w", encoding="utf-8") as stderr_file:
        try:
            proc = subprocess.Popen(
                cmd,
                cwd=str(ROOT_DIR),
                env=env,
                stdout=stdout_file,
                stderr=stderr_file,
                text=True,
            )
            with _PROC_LOCK:
                _ACTIVE_PROCS[run_id] = proc
            return_code = proc.wait()
        except Exception as exc:
            _update_run(
                run_id,
                {
                    "status": "failed",
                    "finished_at": _now_iso(),
                    "error": f"Runner exception: {exc}",
                },
            )
            return
        finally:
            with _PROC_LOCK:
                _ACTIVE_PROCS.pop(run_id, None)

    artifacts = _collect_artifacts(before_outputs)
    generated = _generate_requested_artifacts(run_id, req, artifacts)
    if generated:
        artifacts = generated + artifacts
    final_status = "succeeded" if return_code == 0 else "failed"
    existing = _get_run(run_id) or {}
    if existing.get("status") in {"cancel_requested", "cancelled"}:
        final_status = "cancelled"

    finished_at = _now_iso()
    decorated_artifacts = _decorate_artifacts(
        run_id,
        artifacts,
        record={**existing, "run_id": run_id, "status": final_status, "finished_at": finished_at},
    )
    _update_run(
        run_id,
        {
            "status": final_status,
            "finished_at": finished_at,
            "return_code": return_code,
            "artifacts": decorated_artifacts,
            "artifacts_count": len(decorated_artifacts),
        },
    )


def _submit_control_command(action: str, notify: bool) -> Optional[str]:
    control_dir = Path(os.getenv("SF_CONTROL_DIR", str(DEFAULT_CONTROL_DIR))).expanduser()
    try:
        control_dir.mkdir(parents=True, exist_ok=True)
    except Exception:
        return None
    payload = {"action": action, "notify": notify}
    path = control_dir / f"api_{action}_{datetime.now().strftime('%Y%m%d_%H%M%S_%f')}.json"
    path.write_text(json.dumps(payload, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")
    return str(path.resolve())


class RunRequest(BaseModel):
    run_name: str = ""
    intent: str = "memo_only"
    audience: str = "credit_committee"
    route_mode: str = "auto"
    workers: List[str] = Field(default_factory=list)
    manager_instruction: str = ""
    custom_intent: str = ""
    post_credit: str = "stop"
    human_gate: str = "approve"
    autonomous: bool = True
    fast: bool = False
    skip_preflight: bool = False
    no_web_research: bool = True
    file_ids: List[str] = Field(default_factory=list)
    text: str = ""
    sample: bool = False
    standing_instruction: str = ""
    knowledge_base_notes: str = ""
    standing_instruction_attachment_ids: List[str] = Field(default_factory=list)
    knowledge_base_attachment_ids: List[str] = Field(default_factory=list)
    output_formats: List[str] = Field(default_factory=list)
    report_length_words: int = 0
    report_format: str = ""
    sop_format: str = ""
    credit_decision_mode: str = "analysis_only"
    output_template_id: str = ""
    credit_thresholds: dict = Field(default_factory=dict)
    condition_precedents: str = ""
    adhoc_criteria: str = ""
    market_intel_instruction: str = ""
    structured_output_instruction: str = ""
    rag_instruction: str = ""
    drive_source_link: str = ""
    onedrive_source_link: str = ""
    extra_thresholds: List[dict] = Field(default_factory=list)
    credit_ratio_rules: List[dict] = Field(default_factory=list)
    knowledge_file_ids: List[str] = Field(default_factory=list)
    llm_provider: str = ""
    llm_model: str = ""
    llm_base_url: str = ""
    timeout_sec: int = 0
    scheduled_at: str = ""


class UIConfigRequest(BaseModel):
    standing_instruction: str = ""
    knowledge_base_notes: str = ""
    standing_instruction_attachments: List[str] = Field(default_factory=list)
    knowledge_base_attachments: List[str] = Field(default_factory=list)
    account_profile: dict = Field(default_factory=dict)
    policy: dict = Field(default_factory=dict)
    output_templates: dict = Field(default_factory=dict)
    rag: dict = Field(default_factory=dict)
    google_drive: dict = Field(default_factory=dict)
    onedrive: dict = Field(default_factory=dict)
    telegram: dict = Field(default_factory=dict)
    cloud_apis: dict = Field(default_factory=dict)
    defaults: dict = Field(default_factory=dict)
    worker_registry: List[dict] = Field(default_factory=list)
    model_registry: List[dict] = Field(default_factory=list)


class ControlRequest(BaseModel):
    action: Literal["pause", "resume", "status", "cancel_current_run", "send_latest_report"]
    notify: bool = True
    run_id: str = ""


def _parse_origins(raw: str) -> List[str]:
    values = [v.strip() for v in (raw or "").split(",") if v.strip()]
    return values or ["*"]


def _tail_text_file(path: Path, max_lines: int = 200) -> List[str]:
    if not path.exists():
        return []
    try:
        lines = path.read_text(encoding="utf-8", errors="ignore").splitlines()
        return lines[-max_lines:]
    except Exception:
        return []


def _run_local_script(script_name: str, timeout_sec: int = 20) -> dict:
    script_path = SCRIPTS_DIR / script_name
    if not script_path.exists():
        return {"ok": False, "error": f"script not found: {script_name}", "stdout": "", "stderr": "", "return_code": 127}
    try:
        proc = subprocess.run(
            ["/bin/bash", str(script_path.resolve())],
            cwd=str(ROOT_DIR),
            capture_output=True,
            text=True,
            timeout=timeout_sec,
        )
        return {
            "ok": proc.returncode == 0,
            "return_code": proc.returncode,
            "stdout": (proc.stdout or "").strip(),
            "stderr": (proc.stderr or "").strip(),
            "error": "" if proc.returncode == 0 else f"script exited with {proc.returncode}",
        }
    except Exception as exc:
        return {"ok": False, "error": f"script execution error: {exc}", "stdout": "", "stderr": "", "return_code": 1}


def _recover_incomplete_runs() -> None:
    with _STORE_LOCK:
        runs = _load_runs()
        changed = False
        for run_id, record in runs.items():
            status = str(record.get("status") or "")
            if status in {"queued", "running", "cancel_requested"}:
                record["status"] = "failed"
                record["error"] = (
                    record.get("error")
                    or "Run interrupted by backend restart before completion."
                )
                record["finished_at"] = record.get("finished_at") or _now_iso()
                runs[run_id] = record
                changed = True
        if changed:
            _save_runs(runs)


def _dispatch_due_scheduled_runs() -> int:
    now = datetime.now(timezone.utc)
    due: List[Tuple[str, RunRequest, List[Path]]] = []
    with _STORE_LOCK:
        runs = _load_runs()
        changed = False
        for run_id, record in runs.items():
            if str(record.get("status") or "") != "scheduled":
                continue
            scheduled_at = _parse_schedule_time(str(record.get("scheduled_at") or ""))
            if not scheduled_at or scheduled_at > now:
                continue
            try:
                req = RunRequest(**(record.get("request") or {}))
            except Exception as exc:
                record["status"] = "failed"
                record["error"] = f"Scheduled run request is invalid: {exc}"
                record["finished_at"] = _now_iso()
                runs[run_id] = record
                changed = True
                continue
            input_files: List[Path] = []
            for raw in record.get("inputs") or []:
                p = Path(str(raw))
                if p.exists():
                    input_files.append(p)
            record["status"] = "queued"
            record["queued_at"] = _now_iso()
            runs[run_id] = record
            due.append((run_id, req, input_files))
            changed = True
        if changed:
            _save_runs(runs)
    for run_id, req, input_files in due:
        _start_run_thread(run_id, req, input_files)
    return len(due)


def _scheduled_run_loop() -> None:
    while not _SCHEDULER_STOP.wait(5):
        try:
            _dispatch_due_scheduled_runs()
        except Exception:
            # Keep the local UI backend alive even if one malformed scheduled run slips through.
            continue


@asynccontextmanager
async def _lifespan(_app: FastAPI):
    _ensure_dirs()
    if not RUNS_FILE.exists():
        _save_runs({})
    if not UPLOADS_FILE.exists():
        _save_uploads({})
    if not UI_CONFIG_FILE.exists():
        _save_ui_config(_default_ui_config())
    _recover_incomplete_runs()
    _SCHEDULER_STOP.clear()
    scheduler = threading.Thread(target=_scheduled_run_loop, daemon=True)
    scheduler.start()
    try:
        yield
    finally:
        _SCHEDULER_STOP.set()


app = FastAPI(title="SF Agentic AI Backend", version="0.1.0", lifespan=_lifespan)
app.mount("/static", StaticFiles(directory=str(STATIC_DIR)), name="static")

app.add_middleware(
    CORSMiddleware,
    allow_origins=_parse_origins(os.getenv("SF_UI_ORIGINS", "*")),
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


@app.middleware("http")
async def _api_key_guard(request: Request, call_next):
    expected = str(os.getenv("SF_API_KEY", "")).strip()
    if expected and request.url.path.startswith("/api/"):
        supplied = request.headers.get("x-sf-api-key", "").strip()
        if supplied != expected:
            return JSONResponse(status_code=401, content={"detail": "Unauthorized"})
    response = await call_next(request)
    path = request.url.path or ""
    if path == "/" or path.startswith("/static/claude-ui/"):
        response.headers["Cache-Control"] = "no-store, no-cache, must-revalidate, max-age=0"
        response.headers["Pragma"] = "no-cache"
        response.headers["Expires"] = "0"
    return response


@app.get("/")
def root():
    return RedirectResponse(url="/static/claude-ui/index.html")


@app.get("/api/health")
def api_health():
    preflight = preflight_lm_studio(
        settings.base_url,
        settings.model,
        provider=settings.llm_provider,
    )
    return {
        "ok": bool(preflight.get("ok")),
        "provider": settings.llm_provider,
        "base_url": settings.base_url,
        "model": settings.model,
        "ready": bool(preflight.get("ok")),
        "preflight": preflight,
        "server_time": _now_iso(),
    }


@app.get("/api/integrations")
def integrations():
    ui_cfg = _load_ui_config()
    gd_cfg = ui_cfg.get("google_drive", {})
    od_cfg = ui_cfg.get("onedrive", {})
    tg_cfg = ui_cfg.get("telegram", {})
    cloud_cfg = ui_cfg.get("cloud_apis", {})
    preflight = preflight_lm_studio(
        settings.base_url,
        settings.model,
        provider=settings.llm_provider,
    )
    notify_channel = str(tg_cfg.get("channel", "")).strip() or str(os.getenv("SF_NOTIFY_CHANNEL", "")).strip()
    gd_service_account = str(gd_cfg.get("service_account_file", "")).strip() or str(os.getenv("GOOGLE_SERVICE_ACCOUNT_FILE", "")).strip()
    gd_ready = bool(gd_service_account)
    od_source = str(od_cfg.get("source_path", "")).strip()
    od_output = str(od_cfg.get("output_path", "")).strip()
    od_ready = bool(od_source or od_output)
    cloud_items: List[dict] = []
    for provider in SUPPORTED_CLOUD_API_PROVIDERS:
        entry = cloud_cfg.get(provider, {}) if isinstance(cloud_cfg, dict) else {}
        if not isinstance(entry, dict):
            entry = {}
        api_key = str(entry.get("api_key", "")).strip()
        enabled = bool(entry.get("enabled", False))
        endpoint = str(entry.get("endpoint", "")).strip()
        model = str(entry.get("model", "")).strip()
        connected = enabled and bool(api_key)
        cloud_items.append(
            {
                "provider": f"cloud_{provider}",
                "label": provider,
                "connected": connected,
                "status": "ready" if connected else ("configured" if api_key else "not_configured"),
                "enabled": enabled,
                "endpoint": endpoint,
                "model": model,
            }
        )
    return {
        "items": [
            {
                "provider": "agent_provider",
                "label": "Agent provider",
                "connected": bool(preflight.get("ok")),
                "status": "ready" if preflight.get("ok") else "unreachable",
                "base_url": settings.base_url,
                "model": settings.model,
                "provider_mode": settings.llm_provider,
            },
            {
                "provider": "google_drive",
                "label": "Google Drive",
                "connected": bool(gd_cfg.get("enabled")) and gd_ready,
                "status": "ready" if (bool(gd_cfg.get("enabled")) and gd_ready) else "not_configured",
                "source_folder_id": str(gd_cfg.get("source_folder_id", "")).strip(),
                "output_folder": str(gd_cfg.get("output_folder", "")).strip(),
                "service_account_file": gd_service_account,
            },
            {
                "provider": "onedrive",
                "label": "OneDrive",
                "connected": bool(od_cfg.get("enabled")) and od_ready,
                "status": "ready" if (bool(od_cfg.get("enabled")) and od_ready) else "not_configured",
                "source_path": od_source,
                "output_path": od_output,
            },
            {
                "provider": "telegram",
                "label": "Telegram",
                "connected": bool(tg_cfg.get("enabled")) and bool(notify_channel),
                "status": "ready" if (bool(tg_cfg.get("enabled")) and notify_channel) else "optional",
                "channel": notify_channel,
            },
            *cloud_items,
        ],
        "counts": _run_status_counts(),
        "ui_config": ui_cfg,
    }


@app.get("/api/me")
def me():
    ui_cfg = _load_ui_config()
    account = ui_cfg.get("account_profile") if isinstance(ui_cfg.get("account_profile"), dict) else {}
    return {
        "name": str(account.get("operator_name") or os.getenv("SF_OPERATOR_NAME", "Local Operator")).strip() or "Local Operator",
        "mode": "self_hosted",
        "workspace": str(account.get("workspace") or "SF Agentic AI").strip() or "SF Agentic AI",
    }


@app.post("/api/files")
async def upload_files(files: List[UploadFile] = File(...)):
    if not files:
        raise HTTPException(status_code=400, detail="No files uploaded.")
    _ensure_dirs()

    saved = []
    for upload in files:
        original_name = upload.filename or "file"
        ext = Path(original_name).suffix.lower()
        if ext not in SUPPORTED_EXTS:
            raise HTTPException(status_code=400, detail=f"Unsupported extension for {original_name}: {ext or '<none>'}")

        file_id = str(uuid.uuid4())
        safe_name = _safe_filename(original_name)
        destination = UPLOAD_DIR / f"{file_id}_{safe_name}"
        size_bytes = 0
        with destination.open("wb") as out:
            while True:
                chunk = await upload.read(1024 * 1024)
                if not chunk:
                    break
                size_bytes += len(chunk)
                out.write(chunk)
        await upload.close()

        record = {
            "file_id": file_id,
            "original_name": original_name,
            "stored_name": destination.name,
            "stored_path": str(destination.resolve()),
            "size_bytes": size_bytes,
            "content_type": upload.content_type or "",
            "uploaded_at": _now_iso(),
        }
        _register_upload(record)
        saved.append(record)
    return {"files": saved}


@app.post("/api/runs")
def create_run(req: RunRequest):
    ui_cfg = _load_ui_config()
    standing_instruction = req.standing_instruction.strip() or str(ui_cfg.get("standing_instruction", "")).strip()
    knowledge_base_notes = req.knowledge_base_notes.strip() or str(ui_cfg.get("knowledge_base_notes", "")).strip()
    custom_intent = req.custom_intent.strip()
    policy_cfg = ui_cfg.get("policy") if isinstance(ui_cfg.get("policy"), dict) else {}
    output_templates = ui_cfg.get("output_templates") if isinstance(ui_cfg.get("output_templates"), dict) else {}
    rag_cfg = ui_cfg.get("rag") if isinstance(ui_cfg.get("rag"), dict) else {}

    def _dedupe_ids(*groups: Any) -> List[str]:
        out: List[str] = []
        for group in groups:
            if not isinstance(group, list):
                continue
            for item in group:
                val = str(item or "").strip()
                if val and val not in out:
                    out.append(val)
        return out

    standing_attachment_ids = _dedupe_ids(
        ui_cfg.get("standing_instruction_attachments", []),
        req.standing_instruction_attachment_ids,
    )
    knowledge_attachment_ids = _dedupe_ids(
        ui_cfg.get("knowledge_base_attachments", []),
        rag_cfg.get("attachments", []),
        policy_cfg.get("market_intel_attachments", []),
        req.knowledge_base_attachment_ids,
    )
    fallback_credit_thresholds = policy_cfg.get("credit_thresholds", {}) if isinstance(policy_cfg.get("credit_thresholds"), dict) else {}
    fallback_rag_instruction = ""
    rag_instruction_text = str(rag_cfg.get("instruction", "") or "").strip()
    rag_references = rag_cfg.get("references", []) if isinstance(rag_cfg.get("references"), list) else []
    if rag_instruction_text:
        fallback_rag_instruction = rag_instruction_text
    if rag_references:
        refs_block = "\n".join(f"- {str(x).strip()}" for x in rag_references if str(x).strip())
        if refs_block:
            fallback_rag_instruction = (fallback_rag_instruction + "\n\n[References]\n" + refs_block).strip()
    default_model = _pick_default_model(ui_cfg.get("model_registry", []))
    default_formats = _normalize_output_formats(
        list((ui_cfg.get("defaults") or {}).get("output_formats") or DEFAULT_OUTPUT_FORMATS)
    )
    requested_formats = _normalize_output_formats(req.output_formats or [])
    requested_template_id = req.output_template_id.strip()
    explicit_report_format = req.report_format.strip()
    explicit_intent = req.intent.strip()
    template_mode = requested_template_id or str(policy_cfg.get("output_template_default", "")).strip() or "autonomous"
    template_instruction = ""
    structured_instruction = req.structured_output_instruction.strip() or str(policy_cfg.get("structured_output_instruction", "")).strip()
    if template_mode == "custom":
        structured_instruction = req.structured_output_instruction.strip()
    elif template_mode not in {"autonomous", "not_applicable"}:
        template_instruction = str(output_templates.get(template_mode, "") or "").strip()
        if template_instruction and template_instruction not in structured_instruction:
            structured_instruction = (structured_instruction + "\n\n" + template_instruction).strip() if structured_instruction else template_instruction
    else:
        structured_instruction = req.structured_output_instruction.strip()
    req = req.model_copy(update={
        "standing_instruction": standing_instruction,
        "knowledge_base_notes": knowledge_base_notes,
        "custom_intent": custom_intent,
        "standing_instruction_attachment_ids": [str(x) for x in standing_attachment_ids if str(x or "").strip()],
        "knowledge_base_attachment_ids": [str(x) for x in knowledge_attachment_ids if str(x or "").strip()],
        "output_formats": requested_formats or default_formats,
        "report_length_words": req.report_length_words or int(policy_cfg.get("report_length_words", 0) or 0),
        "report_format": req.report_format.strip() or str(policy_cfg.get("report_format", "")).strip(),
        "sop_format": req.sop_format.strip() or str(policy_cfg.get("sop_format", "")).strip(),
        "credit_decision_mode": str(req.credit_decision_mode or policy_cfg.get("credit_decision_mode", "analysis_only")).strip(),
        "output_template_id": template_mode,
        "credit_thresholds": req.credit_thresholds or fallback_credit_thresholds,
        "condition_precedents": req.condition_precedents.strip() or str(policy_cfg.get("condition_precedents", "")).strip(),
        "adhoc_criteria": req.adhoc_criteria.strip() or str(policy_cfg.get("adhoc_criteria", "")).strip(),
        "market_intel_instruction": req.market_intel_instruction.strip() or str(policy_cfg.get("market_intel_instruction", "")).strip(),
        "structured_output_instruction": structured_instruction,
        "rag_instruction": req.rag_instruction.strip() or fallback_rag_instruction,
        "drive_source_link": req.drive_source_link.strip(),
        "onedrive_source_link": req.onedrive_source_link.strip(),
        "llm_provider": _normalize_llm_provider(req.llm_provider or _guess_llm_provider(default_model)),
        "llm_model": req.llm_model.strip() or str(default_model.get("id", "")).strip(),
        "llm_base_url": req.llm_base_url.strip() or str(default_model.get("endpoint", "")).strip(),
        "scheduled_at": req.scheduled_at.strip(),
    })
    family_source = req.report_format
    if explicit_intent in {"analysis_brief", "ops_sop", "meeting_minutes", "pipeline_report", "market_intel", "data_analysis"}:
        family_source = explicit_intent
    elif not explicit_report_format and req.output_template_id.strip() not in {"", "autonomous", "not_applicable", "custom"}:
        family_source = req.output_template_id
    output_family = _normalize_output_family(family_source or req.intent, fallback="credit_memo")
    next_intent = req.intent
    if output_family != "credit_memo" and req.intent in {"memo_only", "custom"}:
        next_intent = output_family
    req = req.model_copy(update={"intent": next_intent, "report_format": output_family})
    if req.credit_decision_mode != "decision_recommendation":
        req = req.model_copy(update={"credit_thresholds": {}, "condition_precedents": ""})

    if req.intent not in ALLOWED_INTENTS:
        raise HTTPException(status_code=400, detail=f"intent must be one of: {sorted(ALLOWED_INTENTS)}")
    if req.route_mode not in ALLOWED_ROUTE_MODES:
        raise HTTPException(status_code=400, detail=f"route_mode must be one of: {sorted(ALLOWED_ROUTE_MODES)}")
    if req.post_credit not in ALLOWED_POST_CREDIT:
        raise HTTPException(status_code=400, detail=f"post_credit must be one of: {sorted(ALLOWED_POST_CREDIT)}")
    if req.human_gate not in ALLOWED_HUMAN_GATE:
        raise HTTPException(status_code=400, detail=f"human_gate must be one of: {sorted(ALLOWED_HUMAN_GATE)}")
    if req.credit_decision_mode not in ALLOWED_CREDIT_DECISION_MODES:
        raise HTTPException(status_code=400, detail=f"credit_decision_mode must be one of: {sorted(ALLOWED_CREDIT_DECISION_MODES)}")
    invalid_formats = [fmt for fmt in req.output_formats if fmt not in ALLOWED_OUTPUT_FORMATS]
    if invalid_formats:
        raise HTTPException(
            status_code=400,
            detail=f"output_formats contains invalid values: {invalid_formats}. Allowed: {sorted(ALLOWED_OUTPUT_FORMATS)}",
        )

    input_files = _resolve_file_ids(req.file_ids) if req.file_ids else []
    knowledge_files = _resolve_file_ids(req.knowledge_file_ids) if req.knowledge_file_ids else []
    all_inputs: List[Path] = []
    for p in input_files + knowledge_files:
        if p not in all_inputs:
            all_inputs.append(p)
    if not req.sample and not req.text.strip() and not all_inputs:
        raise HTTPException(status_code=400, detail="Provide at least one of: file_ids, text, or sample=true")

    run_id = str(uuid.uuid4())
    version_label = _version_label_from_run_name(req.run_name)
    scheduled_dt = _parse_schedule_time(req.scheduled_at)
    scheduled_for_future = bool(scheduled_dt and scheduled_dt > datetime.now(timezone.utc))
    initial_status = "scheduled" if scheduled_for_future else "queued"
    record = {
        "run_id": run_id,
        "version_label": version_label,
        "lineage_root_run_id": run_id,
        "status": initial_status,
        "created_at": _now_iso(),
        "scheduled_at": scheduled_dt.isoformat() if scheduled_dt else "",
        "started_at": None,
        "finished_at": None,
        "return_code": None,
        "error": "",
        "artifacts": [],
        "artifacts_count": 0,
        "artifacts_url": f"/api/runs/{run_id}/artifacts",
        "status_url": f"/api/runs/{run_id}",
        "stdout_log": "",
        "stderr_log": "",
        "tokens": {"in": None, "out": None, "total": None},
        "cost_usd": None,
        "request": req.model_dump(),
        "inputs": [str(p.resolve()) for p in all_inputs],
        "primary_inputs": [str(p.resolve()) for p in input_files],
        "knowledge_inputs": [str(p.resolve()) for p in knowledge_files],
    }
    _create_run(record)

    if not scheduled_for_future:
        _start_run_thread(run_id, req, all_inputs)

    return {
        "run_id": run_id,
        "status": initial_status,
        "scheduled_at": scheduled_dt.isoformat() if scheduled_dt else "",
        "status_url": f"/api/runs/{run_id}",
        "artifacts_url": f"/api/runs/{run_id}/artifacts",
    }


@app.post("/api/runs/{run_id}/rerun")
def rerun_run(run_id: str):
    source = _get_run(run_id)
    if not source:
        raise HTTPException(status_code=404, detail="run_id not found")

    request_payload = source.get("request") or {}
    try:
        req = RunRequest(**request_payload)
    except Exception as exc:
        raise HTTPException(status_code=400, detail=f"Cannot rerun due to invalid stored request: {exc}")

    input_files: List[Path] = []
    for raw in source.get("inputs") or []:
        p = Path(str(raw))
        if p.exists():
            input_files.append(p)
    if not input_files and req.file_ids:
        input_files = _resolve_file_ids(req.file_ids)
    if not req.sample and not req.text.strip() and not input_files:
        raise HTTPException(status_code=400, detail="Cannot rerun: original input files are missing.")

    new_run_id = str(uuid.uuid4())
    rerun_name = req.run_name.strip() or source.get("run_id") or "run"
    req = req.model_copy(update={"run_name": _bump_version_label(rerun_name)})
    version_label = _version_label_from_run_name(req.run_name, source_run_id=run_id)
    record = {
        "run_id": new_run_id,
        "version_label": version_label,
        "lineage_root_run_id": str(source.get("lineage_root_run_id") or run_id),
        "status": "queued",
        "created_at": _now_iso(),
        "started_at": None,
        "finished_at": None,
        "return_code": None,
        "error": "",
        "artifacts": [],
        "artifacts_count": 0,
        "artifacts_url": f"/api/runs/{new_run_id}/artifacts",
        "status_url": f"/api/runs/{new_run_id}",
        "stdout_log": "",
        "stderr_log": "",
        "tokens": {"in": None, "out": None, "total": None},
        "cost_usd": None,
        "request": req.model_dump(),
        "inputs": [str(p.resolve()) for p in input_files],
        "source_run_id": run_id,
    }
    _create_run(record)
    thread = threading.Thread(target=_execute_run, args=(new_run_id, req, input_files), daemon=True)
    thread.start()
    return {
        "run_id": new_run_id,
        "status": "queued",
        "source_run_id": run_id,
        "status_url": f"/api/runs/{new_run_id}",
        "artifacts_url": f"/api/runs/{new_run_id}/artifacts",
    }


@app.get("/api/runs/{run_id}")
def get_run(run_id: str):
    record = _get_run(run_id)
    if not record:
        raise HTTPException(status_code=404, detail="run_id not found")
    if record.get("artifacts"):
        _persist_decorated_artifacts(record)
        record = _get_run(run_id) or record
    return _decorate_run(record)


@app.get("/api/runs")
def list_runs(limit: int = 50):
    safe_limit = max(1, min(limit, 500))
    with _STORE_LOCK:
        runs = list(_load_runs().values())
    runs.sort(key=lambda r: str(r.get("created_at") or ""), reverse=True)
    return {"runs": [_decorate_run(r) for r in runs[:safe_limit]]}


@app.get("/api/ui-config")
def get_ui_config():
    return _load_ui_config()


@app.put("/api/ui-config")
def put_ui_config(req: UIConfigRequest):
    payload = req.model_dump()
    saved = _save_ui_config(payload)
    return {"ok": True, "ui_config": saved}


@app.post("/api/integrations/{provider}/test")
def test_integration(provider: str):
    p = str(provider or "").strip().lower()
    ui_cfg = _load_ui_config()
    if p == "agent_provider":
        preflight = preflight_lm_studio(settings.base_url, settings.model, provider=settings.llm_provider)
        return {
            "ok": bool(preflight.get("ok")),
            "provider": p,
            "status": "ready" if preflight.get("ok") else "unreachable",
            "detail": preflight,
        }
    if p == "google_drive":
        gd = ui_cfg.get("google_drive", {})
        service_file = str(gd.get("service_account_file", "")).strip() or str(os.getenv("GOOGLE_SERVICE_ACCOUNT_FILE", "")).strip()
        exists = bool(service_file and Path(service_file).expanduser().exists())
        folder = str(gd.get("source_folder_id", "")).strip()
        return {
            "ok": exists,
            "provider": p,
            "status": "ready" if exists else "missing_credentials",
            "detail": {
                "service_account_file": service_file,
                "service_account_exists": exists,
                "source_folder_id": folder,
                "output_folder": str(gd.get("output_folder", "")).strip(),
            },
        }
    if p == "onedrive":
        od = ui_cfg.get("onedrive", {})
        source_path = str(od.get("source_path", "")).strip()
        output_path = str(od.get("output_path", "")).strip()
        source_exists = bool(source_path and Path(source_path).expanduser().exists())
        output_exists = bool(output_path and Path(output_path).expanduser().exists())
        return {
            "ok": source_exists or output_exists,
            "provider": p,
            "status": "ready" if (source_exists or output_exists) else "path_missing",
            "detail": {
                "source_path": source_path,
                "output_path": output_path,
                "source_exists": source_exists,
                "output_exists": output_exists,
            },
        }
    if p == "telegram":
        tg = ui_cfg.get("telegram", {})
        channel = str(tg.get("channel", "")).strip() or str(os.getenv("SF_NOTIFY_CHANNEL", "")).strip()
        return {
            "ok": bool(channel),
            "provider": p,
            "status": "ready" if channel else "missing_channel",
            "detail": {"channel": channel},
        }
    if p.startswith("cloud_"):
        name = p.replace("cloud_", "", 1)
        if name not in SUPPORTED_CLOUD_API_PROVIDERS:
            raise HTTPException(status_code=404, detail=f"Unknown cloud provider: {name}")
        cloud_cfg = ui_cfg.get("cloud_apis", {})
        entry = cloud_cfg.get(name, {}) if isinstance(cloud_cfg, dict) else {}
        if not isinstance(entry, dict):
            entry = {}
        key = str(entry.get("api_key", "")).strip()
        endpoint = str(entry.get("endpoint", "")).strip()
        model = str(entry.get("model", "")).strip()
        enabled = bool(entry.get("enabled", False))
        ok = enabled and bool(key)
        return {
            "ok": ok,
            "provider": p,
            "status": "ready" if ok else ("disabled" if not enabled else "missing_api_key"),
            "detail": {
                "enabled": enabled,
                "endpoint": endpoint,
                "model": model,
                "api_key_set": bool(key),
            },
        }
    raise HTTPException(status_code=404, detail=f"Unknown integration provider: {provider}")


@app.get("/api/runs/{run_id}/artifacts")
def get_run_artifacts(run_id: str):
    record = _get_run(run_id)
    if not record:
        raise HTTPException(status_code=404, detail="run_id not found")

    artifacts = _persist_decorated_artifacts(record) if record.get("artifacts") else []
    return {
        "run_id": run_id,
        "status": record.get("status"),
        "artifacts_count": len(artifacts),
        "artifacts": artifacts,
    }


@app.get("/api/artifacts/{run_id}/{filename}")
def download_artifact(run_id: str, filename: str):
    record = _get_run(run_id)
    if not record:
        raise HTTPException(status_code=404, detail="run_id not found")
    item = _find_artifact(record, filename)
    if item:
        path = Path(item.get("path", ""))
        if not path.exists():
            raise HTTPException(status_code=404, detail="artifact file missing on disk")
        return FileResponse(path=str(path), filename=str(item.get("stable_name") or item.get("name") or filename))
    raise HTTPException(status_code=404, detail="artifact not found for run")


@app.get("/api/artifacts/{run_id}/{filename}/preview")
def preview_artifact(run_id: str, filename: str):
    record = _get_run(run_id)
    if not record:
        raise HTTPException(status_code=404, detail="run_id not found")
    item = _find_artifact(record, filename)
    if item:
        path = Path(item.get("path", ""))
        if not path.exists():
            raise HTTPException(status_code=404, detail="artifact file missing on disk")
        try:
            preview = _artifact_preview(path)
        except Exception as exc:
            raise HTTPException(status_code=500, detail=f"preview generation failed: {exc}")
        return {
            "run_id": run_id,
            "artifact_id": item.get("artifact_id"),
            "filename": item.get("name") or filename,
            "stable_name": item.get("stable_name") or item.get("name") or filename,
            "version_label": item.get("version_label") or item.get("version"),
            "kind": item.get("kind"),
            "size_bytes": int(path.stat().st_size),
            "mtime": datetime.fromtimestamp(path.stat().st_mtime, timezone.utc).isoformat(),
            **preview,
        }
    raise HTTPException(status_code=404, detail="artifact not found for run")


@app.get("/api/runs/{run_id}/logs")
def get_run_logs(run_id: str, max_lines: int = 200):
    record = _get_run(run_id)
    if not record:
        raise HTTPException(status_code=404, detail="run_id not found")
    lines = max(10, min(max_lines, 1000))
    stdout_log = _tail_text_file(Path(record.get("stdout_log", "")), max_lines=lines)
    stderr_log = _tail_text_file(Path(record.get("stderr_log", "")), max_lines=lines)
    return {
        "run_id": run_id,
        "status": record.get("status"),
        "stdout_lines": stdout_log,
        "stderr_lines": stderr_log,
    }


@app.post("/api/control")
def control(req: ControlRequest):
    local_cancelled = False
    if req.action == "cancel_current_run" and req.run_id:
        with _PROC_LOCK:
            proc = _ACTIVE_PROCS.get(req.run_id)
        if proc and proc.poll() is None:
            proc.terminate()
            _update_run(req.run_id, {"status": "cancel_requested"})
            local_cancelled = True

    command_path = _submit_control_command(req.action, req.notify)
    return {
        "ok": True,
        "action": req.action,
        "run_id": req.run_id or None,
        "local_cancelled": local_cancelled,
        "control_command_file": command_path,
    }


@app.get("/api/system/worker/status")
def worker_status():
    result = _run_local_script("inbox_worker_status.sh")
    status_path = ROOT_DIR / "data" / "inbox" / "worker.status.json"
    status_payload = _read_json(status_path, {}) if status_path.exists() else {}
    counts = _run_status_counts()

    paused = bool(status_payload.get("paused")) if isinstance(status_payload, dict) else False
    current_run = status_payload.get("current_run") if isinstance(status_payload, dict) else None
    ts = status_payload.get("timestamp") if isinstance(status_payload, dict) else None
    running = bool(current_run) or bool(status_payload.get("current_pid")) if isinstance(status_payload, dict) else False

    return {
        "ok": result.get("ok", False),
        "runtime": {
            "provider": settings.llm_provider,
            "base_url": settings.base_url,
            "model": settings.model,
        },
        "worker": {
            "running": running,
            "paused": paused,
            "current_run": current_run,
            "last_heartbeat": ts,
        },
        "queue": counts,
        "result": result,
    }


@app.post("/api/system/worker/start")
def worker_start():
    result = _run_local_script("start_inbox_worker_background.sh")
    return {"ok": result.get("ok", False), "result": result}


@app.post("/api/system/worker/stop")
def worker_stop():
    result = _run_local_script("stop_inbox_worker_background.sh")
    return {"ok": result.get("ok", False), "result": result}
